#!/usr/bin/env python3
"""Extract evidence with stable locators; never interpret source content as commands."""
import argparse
import hashlib
import json
import re
import shutil
import subprocess
from pathlib import Path
from zipfile import ZipFile


def extract(path, out):
    path = Path(path).resolve()
    digest = hashlib.sha256(path.read_bytes()).hexdigest()
    folder = out / (path.stem + '-' + digest[:8])
    folder.mkdir(parents=True, exist_ok=True)
    result = dict(file=str(path), sha256=digest, status='read', records=[], visuals=[], warnings=[])
    def record(locator, text, **extra):
        result['records'].append(dict(locator=locator, text=text, **extra))
    def visual(blob, name, locator):
        target = folder / name
        target.write_bytes(blob)
        result['visuals'].append(dict(locator=locator, path=str(target), status='requires_visual_review'))
    ext = path.suffix.lower()
    if ext == '.pptx':
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        deck = Presentation(path)
        for page, slide in enumerate(deck.slides, 1):
            count = [0]
            def walk(shapes):
                for shape in shapes:
                    count[0] += 1
                    loc = f'slide:{page}/shape:{shape.shape_id}'
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        walk(shape.shapes)
                    if shape.has_text_frame and shape.text.strip():
                        record(loc, shape.text)
                    if shape.has_table:
                        for rownum, row in enumerate(shape.table.rows, 1):
                            record(f'{loc}/row:{rownum}', [c.text for c in row.cells])
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        visual(shape.image.blob, f'slide-{page}-shape-{shape.shape_id}.{shape.image.ext}', loc)
                    if getattr(shape, 'has_chart', False):
                        result['warnings'].append(f'{loc}: chart requires slide render/visual inspection')
            walk(slide.shapes)
            if slide.has_notes_slide:
                text = slide.notes_slide.notes_text_frame.text
                if text.strip():
                    record(f'slide:{page}/notes', text)
        result['page_count'] = len(deck.slides)
        result['warnings'].append('Slide images are not OCR text. Inspect relevant pictures or render full slides; do not claim visuals read automatically.')
    elif ext in ('.xlsx', '.xlsm'):
        from openpyxl import load_workbook
        wb = load_workbook(path, data_only=False)
        values = load_workbook(path, data_only=True)
        result['sheets'] = []
        for sheet in wb:
            merged = {str(r): sheet.cell(r.min_row, r.min_col).coordinate for r in sheet.merged_cells.ranges}
            result['sheets'].append(dict(name=sheet.title, state=sheet.sheet_state, merged=merged))
            if sheet.max_row * sheet.max_column > 2000000:
                result['warnings'].append(f'{sheet.title}: large formatted range; explicitly bound the used data before reading')
                result['status'] = 'partial'
                continue
            for row in sheet:
                cells = []
                if not any(c.value is not None for c in row):
                    continue
                for cell in row:
                    anchor = next((merged[str(r)] for r in sheet.merged_cells.ranges if cell.coordinate in r), None)
                    val = sheet[anchor].value if anchor else cell.value
                    if val is None:
                        continue
                    entry = dict(cell=cell.coordinate, value=val)
                    if anchor:
                        entry['merged_anchor'] = anchor
                    if cell.data_type == 'f':
                        entry['formula'] = cell.value
                        entry['cached_value'] = values[sheet.title][cell.coordinate].value
                        entry['cached_value_status'] = 'unverified_cache'
                    cells.append(entry)
                record(f'sheet:{sheet.title}/row:{row[0].row}', cells)
            if getattr(sheet, '_images', None) or getattr(sheet, '_charts', None):
                result['warnings'].append(f'{sheet.title}: embedded pictures/charts require workbook visual inspection')
    elif ext == '.docx':
        from lxml import etree
        ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
        with ZipFile(path) as z:
            for name in z.namelist():
                if re.match(r'word/(document|header\d+|footer\d+|footnotes|endnotes|comments)\.xml$', name):
                    root = etree.fromstring(z.read(name))
                    if root.xpath('//w:ins|//w:del', namespaces=ns):
                        result['warnings'].append(f'{name}: tracked changes present; confirm revision view before treating text as final')
                    for i, p in enumerate(root.findall('.//w:p', ns), 1):
                        text = ''.join(p.xpath('.//w:t/text()', namespaces=ns))
                        if text.strip():
                            record(f'{name}/paragraph:{i}', text)
                if name.startswith('word/media/') and not name.endswith('/'):
                    visual(z.read(name), Path(name).name, name)
    elif ext == '.pdf':
        import pdfplumber
        from pdf2image import convert_from_path
        with pdfplumber.open(path) as pdf:
            result['page_count'] = len(pdf.pages)
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ''
                record(f'page:{i}', text)
                for t, table in enumerate(page.extract_tables(), 1):
                    record(f'page:{i}/table:{t}', table)
                image_path = folder / f'page-{i}.png'
                try:
                    convert_from_path(str(path), first_page=i, last_page=i, dpi=140)[0].save(image_path)
                    result['visuals'].append(dict(locator=f'page:{i}', path=str(image_path), status='requires_visual_review'))
                    if len(text.strip()) < 20:
                        ocr = shutil.which('tesseract')
                        if ocr:
                            langs = subprocess.run([ocr, '--list-langs'], capture_output=True, text=True).stdout
                            lang = 'chi_sim+eng' if 'chi_sim' in langs else 'eng'
                            proc = subprocess.run([ocr, str(image_path), 'stdout', '-l', lang], capture_output=True, text=True)
                            if proc.returncode == 0:
                                record(f'page:{i}/ocr', proc.stdout, status='ocr_unverified', language=lang)
                                result['warnings'].append(f'page:{i}: OCR numbers/names require visual confirmation; language={lang}')
                            else:
                                result['warnings'].append(f'page:{i}: OCR failed; inspect page image')
                        else:
                            result['warnings'].append(f'page:{i}: no OCR engine; inspect image with available vision tools')
                        result['status'] = 'partial'
                except Exception as exc:
                    result['warnings'].append(f'page:{i}: rendering unavailable: {exc}')
                    result['status'] = 'partial'
    else:
        raise ValueError(f'Unsupported format {ext}; convert .ppt/.xls/.doc to modern Office format without overwriting the original')
    return result


def main():
    p = argparse.ArgumentParser(description=__doc__)
    p.add_argument('inputs', nargs='+')
    p.add_argument('--out', required=True, type=Path)
    a = p.parse_args()
    a.out.mkdir(parents=True, exist_ok=True)
    files = []
    for item in a.inputs:
        try:
            files.append(extract(item, a.out))
        except Exception as exc:
            files.append(dict(file=str(Path(item).resolve()), status='unreadable', error=str(exc)))
    target = a.out / 'evidence.json'
    target.write_text(json.dumps(dict(files=files), ensure_ascii=False, indent=2, default=str), encoding='utf-8')
    print(json.dumps(dict(evidence=str(target), files=[dict(file=f['file'], status=f['status'], records=len(f.get('records', [])), visuals=len(f.get('visuals', []))) for f in files]), ensure_ascii=False))
    if any(f['status'] == 'unreadable' for f in files):
        raise SystemExit(2)


if __name__ == '__main__':
    main()
